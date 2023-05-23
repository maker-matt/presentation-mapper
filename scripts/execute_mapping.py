import yaml
import openpyxl
from pptx import Presentation
from utils import get_xlsx_from_dir, get_pptx_from_dir
import os

def get_excel_data(mapping, excel_filepath):
    workbook = openpyxl.load_workbook(excel_filepath)
    data = {}
    def replace_reference_with_value(mapping_value, placeholder):
        # Try to split the mapping_value, if it's a cell reference this will succeed
    
        sheet_name, cell = mapping_value.split('!')

        try:
            worksheet = workbook[sheet_name]
        except KeyError:
            raise KeyError(f"Worksheet '{sheet_name}' not found in workbook '{excel_filepath}'")
        try:
            value = worksheet[cell].value
        except KeyError:
            raise KeyError(f"Cell '{cell}' not found in worksheet '{sheet_name}'")
        if value is None or value == '':
            raise ValueError(f"Cell '{cell}' in worksheet '{sheet_name}' is empty")
        else:
            data[placeholder] = value

    for placeholder, mapping_value in mapping.items():
        try:
            replace_reference_with_value(mapping_value, placeholder)
        except ValueError:
            # If it's not a cell reference, ValueError will be raised and we can treat it as a direct value, and leave it as-is
            continue
        except AttributeError:
            raise AttributeError(f"Invalid mapping value '{mapping_value}' for placeholder '{placeholder}'. Did you forget to use \"quotes\" in mapping.yml?")
    workbook.close()
    return data

def read_mapping(mapping_filepath):
    with open(mapping_filepath, 'r') as f:
        mapping = yaml.safe_load(f)
    return flatten_mapping(mapping)

def flatten_mapping(mapping):
    flattened_mapping = {}
    for slide, placeholders in mapping.items():
        flattened_mapping.update(placeholders)
    return flattened_mapping

def replace_text_with_mapping(text, mapping):
    for key, value in mapping.items():
        text = text.replace('{{' + key + '}}', str(value))
    return text

def replace_placeholders(presentation, mapping):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = replace_text_with_mapping(run.text, mapping)
    return presentation

def execute_mapping(presentation_filepath, mapping_filepath, excel_filename):

    # Get the required objects
    mapping = read_mapping(mapping_filepath)
    excel_data = get_excel_data(mapping, excel_filename)

    # Combine the mapping with the data from Excel
    mapping.update(excel_data)

    # Replace placeholders in the presentation
    presentation = Presentation(presentation_filepath)
    presentation = replace_placeholders(presentation, mapping)

    # just presentation filename 
    output_filename = os.path.basename(presentation_filepath)[:-len(".pptx")] + '_updated.pptx'
    # change parent directory to "output_files"
    output_filepath = os.path.join(os.path.dirname(presentation_filepath), '..', 'output_files', output_filename)
    presentation.save(output_filepath)

    print(f"Updated presentation '{presentation_filepath}', saved at '{output_filepath}'")

def get_filename(type_of_file, filetype):
    filename = input(f"Enter the {type_of_file} filename: (include the {filetype} extension)")
    if not filename.endswith(f".{filetype}"):
        raise ValueError(f"Invalid filename '{filename}'")
    return filename

if __name__ == '__main__':
    presentation_filename = get_pptx_from_dir(os.path.join(os.getcwd(), 'source_files'))
    excel_filename = get_xlsx_from_dir(os.path.join(os.getcwd(),'source_files'))
    execute_mapping(presentation_filename, os.path.join(os.getcwd(), 'source_files', 'mapping.yml'), excel_filename)
