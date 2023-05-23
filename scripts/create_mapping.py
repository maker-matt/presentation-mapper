
from pptx import Presentation
from utils import extract_keywords, save_mapping, get_pptx_from_dir
import os

def generate_mapping(presentation_path):
    # Load the PowerPoint presentation
    prs = Presentation(presentation_path)

    mapping = {}
    used_keywords = set()

    # Iterate over the slides and extract keywords
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_keywords = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    extracted_keywords = extract_keywords(text)
                    slide_keywords.extend(extracted_keywords)

        # Add unique keywords from the slide to the mapping
        for keyword in set(slide_keywords):
            if keyword not in used_keywords:
                slide_key = f'slide{slide_num}'
                mapping.setdefault(slide_key, {})[keyword] = None
                used_keywords.add(keyword)

    return mapping

if __name__ == '__main__':
    mapping = generate_mapping(get_pptx_from_dir(os.path.join(os.getcwd(),'source_files')))
    save_mapping(mapping, os.path.join(os.getcwd(), 'source_files', 'mapping.yml'))